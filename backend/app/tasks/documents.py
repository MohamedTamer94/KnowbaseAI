from sqlalchemy.orm.session import Session

from app.crud.document_crud import update_document_status
from app.models.document import Document, DocumentType
from app.tasks.crawler import crawl_website
from app.crud.chat_crud import load_recent_messages

from sentence_transformers import SentenceTransformer

import fitz

from app.models.chat import ChatMessage
from docx import Document as DocxDocument
from pptx import Presentation
from bs4 import BeautifulSoup

import chromadb

from app.config import settings 
from app.utils.groq_client import chat_completion_stream, rewrite_query_if_needed

from pathlib import Path

model = SentenceTransformer(settings.SENTENCE_TRANSFORMER_MODEL)

chroma_client = chromadb.HttpClient(host=settings.CHROMA_HOST, port=settings.CHROMA_PORT)

def process_document(db: Session, document: Document):
    try:
        update_document_status(db, document, "processing")

        # Read the document and extract text according to its mime type
        text_blocks = extract_text(document)
        # cheunk the text
        chunks = chunk_text(text_blocks)
        # Embedding
        vectors = embed_chunks(chunks)

        # Store vectors in ChromaDB
        store_vectors_in_chroma(vectors, chunks, document)

        update_document_status(db, document, "done")
    except Exception as e:
        update_document_status(db, document, "failed")
        print(f"Error processing document {document.id}: {e}")

def process_website(db: Session, document: Document, url: str):
    try:
        update_document_status(db, document, "processing")

        print("Starting website crawl for URL:", url)
        blocks = crawl_website(url, max_depth=1)
        print("Crawled blocks:", blocks)
        chunks = chunk_text(blocks)
        embeddings = embed_chunks(chunks)

        store_vectors_in_chroma(embeddings, chunks, document)

        update_document_status(db, document, "done")

    except Exception as e:
        update_document_status(db, document, "failed")
        print(f"Website crawl failed: {e}")

def extract_text(document: Document):
    mimetype = document.type
    path = Path(document.filepath)

    if mimetype == DocumentType.pdf:
        return extract_pdf(path)
    elif mimetype == DocumentType.docx:
        return extract_docx(path)
    elif mimetype == DocumentType.pptx:
        return extract_pptx(path)
    elif mimetype == DocumentType.text:
        return extract_text_file(path)
    elif mimetype == DocumentType.website:
        return extract_html(path)
    else:
        raise ValueError(f"Unsupported MIME type: {mimetype}")


def extract_pdf(path: Path):
    doc = fitz.open(path)
    blocks = []

    for page in doc:
        text = page.get_text().strip()
        if text:
            blocks.append({
                "text": text,
                "source": {"page": page.number + 1}
            })

    return blocks

def extract_docx(path: Path):
    doc = DocxDocument(path)
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    return [{"text": text, "source": {}}]

def extract_pptx(path: Path):
    prs = Presentation(path)
    blocks = []

    for slide_number, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text.strip())
        text = "\n".join(slide_text).strip()
        if text:
            blocks.append({
                "text": text,
                "source": {"slide": slide_number + 1}
            })
    return blocks

def extract_text_file(path: Path):
    text = path.read_text(encoding="utf-8", errors="ignore")
    return [{"text": text, "source": {}}]

def extract_html(path: Path):
    html = path.read_text(encoding="utf-8", errors="ignore")
    soup = BeautifulSoup(html, "html.parser")
    
    for tag in soup(["script", "style", "nav", "footer"]):
        tag.decompose()

    text = soup.get_text(separator="\n")
    return [{"text": text.strip(), "source": {}}]

def chunk_text(blocks, chunk_size=250, overlap=50):
    chunks = []

    for block in blocks:
        words = block["text"].split()
        start = 0

        while start < len(words):
            end = start + chunk_size
            chunk_text = " ".join(words[start:end])

            chunks.append({
                "text": chunk_text,
                "source": block.get("source", {})
            })

            start += chunk_size - overlap

    return chunks

def embed_chunks(chunks: list[dict]):
    return model.encode([chunk["text"] for chunk in chunks], show_progress_bar=True)

def get_collection(tenant_id: str):
    return chroma_client.get_or_create_collection(name=tenant_id, metadata={"hnsw:space": "cosine"})

def store_vectors_in_chroma(vectors, chunks, document: Document):
    collection = get_collection(document.tenant_id)

    ids = [f"{document.id}_{i}" for i in range(len(chunks))]
    metadatas = [
        {
            "doc_id": str(document.id),
            "chunk_id": i,
            "doc_name": document.filename,
            "source": json.dumps(chunk["source"])
        }
        for i, chunk in enumerate(chunks)
    ]

    collection.add(
        ids=ids,
        embeddings=vectors,
        documents=[chunk["text"] for chunk in chunks],
        metadatas=metadatas
    )

def build_prompt(query: str, chunks: list[str]) -> str:
    prompt = "Answer using only the context below. If you don't know, say so.\n\n"
    for chunk in chunks:
        prompt += f"- {chunk}\n"
    prompt += f"\nQuestion: {query}\nAnswer:"
    return prompt


import json

def query_document_stream(db, tenant_id, session_id, query):
    history = load_recent_messages(db, session_id)
    rewritten_query = rewrite_query_if_needed(query, history)
    print("Rewritten query:", rewritten_query)
    collection = get_collection(tenant_id)
    query_embedding = embed_chunks([{"text": rewritten_query}])

    results = collection.query(
        query_embeddings=query_embedding,
        n_results=5
    )

    documents = []
    sources = []

    for doc, meta, score in zip(
        results['documents'][0],
        results['metadatas'][0],
        results['distances'][0]
    ):
        print(meta)
        if score >= 0.2:
            documents.append(doc)
            source = json.loads(meta["source"])
            db_doc = db.query(Document).filter(Document.id == meta['doc_id']).one_or_none()
            sources.append({
                "document_id": meta['doc_id'],
                "doc_name": db_doc.filename if db_doc else "Unknown",
                "source": source,
                "chunk_id": meta['chunk_id'],
                "text": doc[:400]
            })

    def stream():
        # start event
        yield json.dumps({"type": "start", "session_id": str(session_id)}) + "\n"

        # token events
        assistant_text = ""
        for token in ask_llm_stream(query, documents, history):
            assistant_text += token
            yield json.dumps({
                "type": "token",
                "data": token
            }) + "\n"

        # sources event
        yield json.dumps({
            "type": "sources",
            "data": sources
        }) + "\n"
    
         # persist messages
        db.add(ChatMessage(
            session_id=session_id,
            role="user",
            content=query
        ))
        db.add(ChatMessage(
            session_id=session_id,
            role="assistant",
            content=assistant_text
        ))
        db.commit()

    return stream()

def ask_llm_stream(query, chunks, history):
    prompt = build_prompt(query, chunks)

    messages = [
        {"role": "system", "content": "You are a helpful assistant."}
    ]
    for msg in history:
        messages.append({
            "role": msg.role,
            "content": msg.content
        })
    messages.append({
        "role": "user",
        "content": prompt
    })

    buffer = ""

    for token in chat_completion_stream(
        messages=messages,
        temperature=0.2
    ):
        buffer += token

        # flush every ~20 chars for smoother UX
        if len(buffer) >= 20:
            yield buffer
            buffer = ""

    if buffer:
        yield buffer

def delete_document_vectors(document: Document):
    collection = get_collection(document.tenant_id)
    collection.delete(
        where={"doc_id": str(document.id)}
    )